from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path


OUT = Path(__file__).resolve().parent / "OTC_论文复现汇报_模板版.pptx"


def set_bg(slide, color=(245, 248, 252)):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(12.0), Inches(0.9))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(34)
    r.font.bold = True
    r.font.color.rgb = RGBColor(21, 45, 90)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.72), Inches(1.2), Inches(12.0), Inches(0.5))
        st = sub.text_frame
        st.clear()
        sp = st.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.name = "Microsoft YaHei"
        sr.font.size = Pt(16)
        sr.font.color.rgb = RGBColor(70, 88, 130)


def add_bullets(slide, x, y, w, h, lines, title=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    if title:
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.name = "Microsoft YaHei"
        r.font.size = Pt(20)
        r.font.bold = True
        r.font.color.rgb = RGBColor(26, 57, 109)
    for i, t in enumerate(lines):
        p = tf.add_paragraph() if (title or i > 0) else tf.paragraphs[0]
        p.text = t
        p.level = 0
        p.space_before = Pt(6)
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(35, 47, 62)


def add_note_box(slide, x, y, w, h, text, title="讲解提示"):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(234, 241, 255)
    shape.line.color.rgb = RGBColor(145, 171, 217)
    tf = shape.text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = f"{title}："
    r0.font.name = "Microsoft YaHei"
    r0.font.size = Pt(15)
    r0.font.bold = True
    r0.font.color.rgb = RGBColor(41, 73, 132)
    p1 = tf.add_paragraph()
    p1.text = text
    p1.font.name = "Microsoft YaHei"
    p1.font.size = Pt(14)
    p1.font.color.rgb = RGBColor(45, 61, 90)


def add_section_divider(prs, idx, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, (231, 238, 252))
    add_title(slide, f"Part {idx}. {title}", subtitle)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, (233, 240, 252))
    add_title(s, "OTC论文复现与优化汇报", "Optimal Transport Enhanced Cross-City Site Recommendation")
    add_bullets(
        s, 0.8, 2.1, 11.8, 2.6,
        [
            "汇报人：XXX    指导教师：XXX",
            "数据集：OpenSiteRec（Chicago / NYC / Singapore / Tokyo）",
            "目标：完整复现 + 差距分析 + 改进提升",
        ],
    )
    add_note_box(s, 0.8, 5.3, 11.8, 1.5, "开场30秒先说清楚‘复现是否完全成功’：趋势复现 + 数值仍有差距。")

    # 2. Agenda
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "目录")
    add_bullets(
        s, 1.0, 1.5, 11.2, 4.8,
        [
            "1) 问题分析：论文主题、模型结构与数学意义",
            "2) 实验设置：论文协议、代码实现、关键修复",
            "3) 结果与差距分析：多图表对比 + 原因解释",
            "4) 改进提升：在复现基础上的可落地方向",
        ],
    )

    # Part 1 divider
    add_section_divider(prs, "1", "问题分析", "从问题定义到公式意义，建立统一理解")

    # 3. Problem definition
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "1.1 论文问题定义")
    add_bullets(
        s, 0.7, 1.4, 6.0, 4.8,
        [
            "任务：给定品牌（Brand），在目标城市推荐Top-N区域（Region）。",
            "难点：单城市数据稀疏且长尾明显，冷启动品牌/区域多。",
            "核心思路：引入其他城市知识，但避免负迁移。",
        ],
        title="Cross-City Site Recommendation"
    )
    add_note_box(s, 7.0, 1.4, 5.6, 4.8, "这里建议配一张‘四城市知识迁移示意图’（可后续替换图片）。", "图示占位")

    # 4. Model & formula
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "1.2 模型结构与关键公式")
    add_bullets(
        s, 0.7, 1.4, 6.3, 5.0,
        [
            "Base模型：VanillaMF / LightGCN，先学习单城品牌-区域偏好。",
            "OTC：对源城与目标城嵌入做GW对齐，得到推断评分。",
            "最终评分：目标城原始分数 + 多源城推断分数加权融合。",
        ],
    )
    add_note_box(
        s, 7.2, 1.4, 5.3, 5.0,
        "建议在此页放两条公式：\n1) Base评分函数\n2) OTC融合公式（含γ）\n并用一句话解释：γ越大，迁移信息影响越强。",
        "公式讲解框"
    )

    # 5. Formula meaning
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "1.3 数学公式的意义解析")
    add_bullets(
        s, 0.8, 1.4, 11.7, 5.3,
        [
            "GW对齐本质：让“关系结构相似”的品牌/区域在跨城市中匹配，而非逐点硬对齐。",
            "融合权重γ：平衡“目标城自有知识”与“外部迁移知识”；γ过大容易负迁移。",
            "per-source γ：每个源城单独赋权，可抑制不相似城市干扰。",
            "这解释了为什么同一OTC在不同目标城市收益差异很大。",
        ],
    )

    # Part 2 divider
    add_section_divider(prs, "2", "实验设置", "最重要：论文协议 ↔ 代码实现 ↔ 复现质量")

    # 6. paper settings
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "2.1 论文实验协议提炼")
    add_bullets(
        s, 0.7, 1.35, 6.2, 5.2,
        [
            "数据：OpenSiteRec四城市；品牌5-core筛选。",
            "划分：每个品牌内 70% / 10% / 20%（train / val / test）。",
            "指标：Recall@20, nDCG@20。",
            "超参：lr、batch、γ网格搜索，early stopping选最优epoch。",
        ],
    )
    add_note_box(s, 7.1, 1.35, 5.4, 5.2, "此页建议插入论文Table2/4.1.4中的关键句截图。", "证据材料")

    # 7. your setting
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "2.2 我的实现设置（论文未明确参数）")
    add_bullets(
        s, 0.8, 1.45, 11.6, 4.9,
        [
            "固定随机种子 + split_seed分离控制，支持可重复重建数据划分。",
            "LightGCN扩展可调层数（2/3/4）、嵌入维度（64/100/128）。",
            "OTC采用per-source gamma坐标搜索，且允许zero-gamma抑制负迁移。",
            "额外保存中间指标JSON与embedding，便于故障定位与复现实验流水线。",
        ],
    )

    # 8. code framework
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "2.3 代码框架与实验流水线")
    add_bullets(
        s, 0.7, 1.35, 6.3, 5.3,
        [
            "data_utils.py：5-core过滤 + 按品牌分层划分 + 图构建。",
            "main.py：单城训练、评估、早停、导出embedding。",
            "otc_lightgcn.py：GW对齐、分数融合、γ搜索、全城市评估。",
            "run_*.py：自动化调参、批量实验、结果汇总。",
        ],
    )
    add_note_box(
        s, 7.2, 1.35, 5.3, 5.3,
        "可画流程图：\n数据预处理→单城训练→保存embedding→OTC融合→指标统计。\n这页是老师最看重的“你是否真的做了工程复现”。",
        "流程图占位"
    )

    # 9. issues & fixes
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "2.4 复现过程中的关键问题与修复")
    add_bullets(
        s, 0.8, 1.35, 11.6, 5.6,
        [
            "问题1：数据划分不可控，导致结果波动大。修复：force_split + split_seed。",
            "问题2：OTC统一γ易导致新加坡负迁移。修复：per-source γ + zero-gamma。",
            "问题3：OTC选参目标不匹配论文重点。修复：改为nDCG优先。",
            "问题4：长流程难定位误差。修复：指标JSON落盘、分阶段脚本化。",
            "结论：代码层修复比单纯多跑seed更能提升复现质量。",
        ],
    )

    # Part 3 divider
    add_section_divider(prs, "3", "结果与差距分析", "多维结果可视化 + 可解释的差距归因")

    # 10. result table
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "3.1 主结果总览（论文 vs 复现）")
    add_note_box(
        s, 0.8, 1.4, 11.8, 4.8,
        "请在此插入总表（建议8行）：\nChicago/NYC/Singapore/Tokyo × (LightGCN, OTC-LightGCN)\n列：Paper Rec@20 / Repro Rec@20 / Gap / Paper nDCG@20 / Repro nDCG@20 / Gap",
        "总表占位（必须）"
    )
    add_bullets(s, 0.8, 6.0, 11.8, 1.0, ["讲解顺序建议：先看趋势一致性，再看绝对值差距。"])

    # 11. chart page
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "3.2 图表全景分析")
    add_bullets(
        s, 0.8, 1.4, 11.8, 5.7,
        [
            "图1：四城市 Recall@20 柱状对比（Paper / Repro / OTC-Repro）。",
            "图2：四城市 nDCG@20 柱状对比（同上）。",
            "图3：每城市 OTC 相对提升率（%）对比图。",
            "图4：关键超参（γ / layer / split_seed）敏感性折线图。",
            "图5：Singapore 专项（负迁移→抑制后）前后对比图。",
        ],
    )

    # 12. gap explanation
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "3.3 同数据集仍有差距的原因解释")
    add_bullets(
        s, 0.8, 1.4, 11.8, 5.6,
        [
            "数据口径差异：字段映射、候选集定义、去重策略会改变评估难度。",
            "协议细节差异：split随机性、early-stop标准、未公开参数默认值。",
            "OT数值优化差异：GW/entropic-GW实现与收敛设置影响迁移质量。",
            "工程实现差异：embedding导出层、融合顺序、归一化策略等。",
            "结论：复现不等于复刻，需把‘差距来源’证据化展示。",
        ],
    )

    # Part 4 divider
    add_section_divider(prs, "4", "改进提升", "在复现基础上做增量创新（老师最关心）")

    # 13. improvement ideas
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "4.1 可落地改进方向（建议选择其一主打）")
    add_bullets(
        s, 0.8, 1.4, 11.8, 5.8,
        [
            "方向A（推荐）：源城市自适应门控（source gating）替代固定γ网格。",
            "方向B：基于城市相似度先验的γ初始化，再做局部搜索。",
            "方向C：多目标优化（Recall+nDCG）而非单目标选参。",
            "方向D：稳健OT（更强正则+温度缩放）降低数值不稳定与负迁移。",
            "方向E：跨城市对齐前先做语义层（品牌/类别）预聚类。",
        ],
    )

    # 14. your concrete plan
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_title(s, "4.2 本次建议主改进：自适应源城市门控OTC")
    add_bullets(
        s, 0.8, 1.4, 11.8, 5.8,
        [
            "核心：学习每个源城市对目标城市的贡献权重（替代手工γ）。",
            "实现：以验证集nDCG为目标，优化可微门控参数g_s∈[0,1]。",
            "预期收益：减少负迁移、提升跨城市泛化稳定性。",
            "实验设计：与固定γ、per-source γ做消融对比，报告均值±方差。",
            "风险与控制：防止过拟合（L1稀疏约束 + 早停 + 多seed验证）。",
        ],
    )

    # 15. conclusion
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, (233, 240, 252))
    add_title(s, "总结与答辩准备")
    add_bullets(
        s, 0.9, 1.5, 11.6, 4.6,
        [
            "复现价值：复现实验不仅是“跑出数值”，更是“定位差距来源”。",
            "当前结论：趋势已基本可复现，数值差距可解释且可继续缩小。",
            "改进方向：已给出可实现、可评估、可答辩的增量方案。",
            "下一步：完成最终重跑、更新图表、提交论文式复现报告。",
        ],
    )
    add_note_box(s, 0.9, 6.0, 11.6, 1.0, "答辩时最后一句建议：‘我不仅复现了结果，也复现了论文背后的假设与边界。’")

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
